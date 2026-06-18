"""Text extraction, one extractor per supported file type.

Each extractor returns plain text. Keeping them in a single dispatch table makes
adding a new file type a one-line change, and keeps file-format dependencies
isolated from the rest of the pipeline.
"""
from __future__ import annotations

from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

# Maps a normalized document_type to the extractor function.
SUPPORTED_TYPES = ("pdf", "pptx", "docx", "txt", "md")


def detect_type(filename: str) -> str | None:
    ext = Path(filename).suffix.lower().lstrip(".")
    if ext == "markdown":
        ext = "md"
    return ext if ext in SUPPORTED_TYPES else None


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                parts.append(shape.text_frame.text)
    return "\n\n".join(parts)


def _extract_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


_EXTRACTORS = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "txt": _extract_text,
    "md": _extract_text,
}


def extract_text(path: Path, document_type: str) -> str:
    extractor = _EXTRACTORS.get(document_type)
    if extractor is None:
        raise ValueError(f"Unsupported document_type: {document_type}")
    return extractor(path).strip()
